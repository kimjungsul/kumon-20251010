import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLUE_DARK  = RGBColor(0x1F, 0x35, 0x64)
BLUE_MID   = RGBColor(0x2E, 0x75, 0xB6)
BLUE_LIGHT = RGBColor(0xD6, 0xE4, 0xF0)
GREEN      = RGBColor(0x37, 0x86, 0x44)
GREEN_LIGHT= RGBColor(0xE2, 0xEF, 0xDA)
ORANGE     = RGBColor(0xED, 0x7D, 0x31)
ORANGE_LIGHT=RGBColor(0xFC, 0xE4, 0xD6)
RED        = RGBColor(0xC0, 0x00, 0x00)
RED_LIGHT  = RGBColor(0xF4, 0xCC, 0xCC)
PURPLE     = RGBColor(0x70, 0x30, 0xA0)
PURPLE_LIGHT=RGBColor(0xED, 0xE7, 0xF6)
YELLOW     = RGBColor(0xFF, 0xC0, 0x00)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_LIGHT = RGBColor(0xF5, 0xF5, 0xF5)
GRAY_MID   = RGBColor(0xD9, 0xD9, 0xD9)
BLACK      = RGBColor(0x00, 0x00, 0x00)

blank = prs.slide_layouts[6]

def add_rect(slide, l, t, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    return shape

def add_text(slide, text, l, t, w, h, size=16, bold=False, color=None,
             align=PP_ALIGN.LEFT, wrap=True):
    if color is None:
        color = BLACK
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox

def add_multiline(slide, lines, l, t, w, h, size=14, color=None, bold_first=False):
    """lines: list of (text, bold) tuples"""
    if color is None:
        color = BLACK
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for (text, bold) in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return txBox

# ────────────────────────────────────────────────────────────
# 슬라이드 1: 표지
# ────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
add_rect(slide, 0, 0, 13.33, 7.5, BLUE_DARK)
add_rect(slide, 0, 5.6, 13.33, 1.9, BLUE_MID)
add_rect(slide, 1.5, 1.0, 10.33, 0.08, WHITE)

add_text(slide, '스마트구몬N', 0, 1.5, 13.33, 1.2,
         size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, '유아 회원 설득 화법', 0, 2.8, 13.33, 0.9,
         size=30, color=RGBColor(0xBD, 0xD7, 0xEE), align=PP_ALIGN.CENTER)
add_text(slide, '학습 과목 유형별 상담 가이드', 0, 3.65, 13.33, 0.6,
         size=18, color=RGBColor(0x9D, 0xC3, 0xE6), align=PP_ALIGN.CENTER)

# 특징 3개
features = ['종이 + 앱 병행', '어떤 패드든 설치 가능', '월 +23,000원 · N페이 1만원']
feat_colors = [ORANGE, GREEN, YELLOW]
for i, (ft, fc) in enumerate(zip(features, feat_colors)):
    x = 1.2 + i * 3.7
    add_rect(slide, x, 4.55, 3.2, 0.75, fc)
    add_text(slide, ft, x, 4.62, 3.2, 0.6,
             size=15, bold=True, color=BLUE_DARK, align=PP_ALIGN.CENTER)

add_text(slide, '구몬 관양지국  |  교사 상담 자료  |  대상: 4~7세 유아', 0, 6.2, 13.33, 0.5,
         size=13, color=WHITE, align=PP_ALIGN.CENTER)

# ────────────────────────────────────────────────────────────
# 슬라이드 2: 4가지 유형 개요
# ────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
add_rect(slide, 0, 0, 13.33, 1.05, BLUE_DARK)
add_text(slide, '학습 과목 유형별 분류  —  총 10명', 0.3, 0.15, 12, 0.75,
         size=24, bold=True, color=WHITE)

types = [
    (BLUE_MID,  BLUE_LIGHT,  '유형 1',  '수학 + 국어\n기본형',
     '허예성  수학2A · 국어3A\n손채희  수학2A · 국어2A\n정은기  수학A · 국어3A',
     '두 과목 복습을 앱으로 연결\n→ 학습 완성도 강화'),
    (GREEN,     GREEN_LIGHT, '유형 2',  '다과목형\n(3~4과목)',
     '정하윤  수학2A · 국어3A · 한자5A · SSR\n서아준  수학A · 국어4A · 한자5A\n박시환  수학2A · 국어3A · 한자5A\n이은서  수학2A · 국어2A · 한자5A · 호과2',
     '과목 추가 아님\n→ 기존 학습을 앱으로 효율화'),
    (ORANGE,    ORANGE_LIGHT,'유형 3',  '운필·쓰기\n훈련 중심',
     '채아린  수학2A · 국어3A · 운필4\n채아성  국어3A · 한자4A · 운필7',
     '종이 손 훈련 + 앱 흥미 연결\n→ 가장 시너지 높은 유형'),
    (PURPLE,    PURPLE_LIGHT,'유형 4',  '국어 단독',
     '임윤아  국어A',
     '앱으로 흥미 붙여\n→ 수학 과목 추가까지 연결'),
]

for i, (hc, bc, badge, title, members, point) in enumerate(types):
    x = 0.25 + i * 3.27
    add_rect(slide, x, 1.2, 3.0, 5.85, bc)
    add_rect(slide, x, 1.2, 3.0, 0.5, hc)
    add_text(slide, badge, x + 0.05, 1.27, 2.9, 0.38,
             size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, title, x + 0.1, 1.8, 2.8, 0.75,
             size=14, bold=True, color=hc)
    add_text(slide, members, x + 0.1, 2.65, 2.8, 2.4, size=11, color=BLACK)
    add_rect(slide, x + 0.1, 5.2, 2.8, 0.04, hc)
    add_text(slide, '상담 포인트', x + 0.1, 5.3, 2.8, 0.3,
             size=11, bold=True, color=hc)
    add_text(slide, point, x + 0.1, 5.65, 2.8, 1.2, size=11, color=BLACK)

# ────────────────────────────────────────────────────────────
# 슬라이드 3: 유형 1 — 수학+국어 기본형
# ────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
add_rect(slide, 0, 0, 13.33, 1.05, BLUE_MID)
add_rect(slide, 0, 0, 0.55, 1.05, BLUE_DARK)
add_text(slide, '유형 1', 0.02, 0.28, 0.52, 0.5,
         size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, '수학 + 국어 기본형  |  허예성 · 손채희 · 정은기', 0.65, 0.18, 12, 0.7,
         size=22, bold=True, color=WHITE)

# 왼쪽: 회원 정보
add_rect(slide, 0.25, 1.2, 4.0, 5.8, GRAY_LIGHT)
add_rect(slide, 0.25, 1.2, 4.0, 0.45, BLUE_MID)
add_text(slide, '대상 회원', 0.4, 1.27, 3.7, 0.32, size=13, bold=True, color=WHITE)

members_info = [
    ('허예성', '수학2A · 국어3A'),
    ('손채희', '수학2A · 국어2A'),
    ('정은기', '수학A · 국어3A'),
]
for i, (name, subj) in enumerate(members_info):
    y = 1.85 + i * 1.65
    add_rect(slide, 0.35, y, 3.8, 1.45, WHITE)
    add_rect(slide, 0.35, y, 0.08, 1.45, BLUE_MID)
    add_text(slide, name, 0.55, y + 0.1, 3.4, 0.45, size=15, bold=True, color=BLUE_DARK)
    add_text(slide, subj, 0.55, y + 0.6, 3.4, 0.6, size=12, color=BLACK)

# 오른쪽: 화법
add_rect(slide, 4.55, 1.2, 8.55, 5.8, GRAY_LIGHT)
add_rect(slide, 4.55, 1.2, 8.55, 0.45, BLUE_DARK)
add_text(slide, '상담 화법', 4.7, 1.27, 8.2, 0.32, size=13, bold=True, color=WHITE)

sections = [
    (BLUE_MID, '도입 (공감)',
     '"저도 이 나이에 억지로 공부시키는 건 절대 안 권해요.\n스마트구몬N은 공부가 아니라 놀이처럼 시작하는 습관이에요."'),
    (BLUE_DARK, '전환 (차별점)',
     '"예성이(채희/은기)도 지금 수학이랑 국어 두 과목 하고 있잖아요.\n종이 교재로 집중력을 키우고, 집에서는 갖고 있는 패드에 앱 깔아서 복습해요.\n어떤 기기든 설치되니까 따로 뭘 살 필요도 없어요."'),
    (GREEN, '클로징',
     '"지금 이 시기에 앉아서 뭔가 해보는 경험이 초등 가서 엄청난 차이를 만들어요.\n한번 체험만 해보시겠어요?"'),
]
for i, (sc, label, text) in enumerate(sections):
    y = 1.85 + i * 1.7
    add_rect(slide, 4.65, y, 8.35, 1.5, WHITE)
    add_rect(slide, 4.65, y, 0.08, 1.5, sc)
    add_text(slide, label, 4.85, y + 0.05, 7.9, 0.35, size=11, bold=True, color=sc)
    add_text(slide, text, 4.85, y + 0.42, 7.9, 1.0, size=12, color=BLACK)

# ────────────────────────────────────────────────────────────
# 슬라이드 4: 유형 2 — 다과목형
# ────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
add_rect(slide, 0, 0, 13.33, 1.05, GREEN)
add_rect(slide, 0, 0, 0.55, 1.05, RGBColor(0x1E, 0x5C, 0x2E))
add_text(slide, '유형 2', 0.02, 0.28, 0.52, 0.5,
         size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, '다과목형 (3~4과목)  |  정하윤 · 서아준 · 박시환 · 이은서', 0.65, 0.18, 12, 0.7,
         size=22, bold=True, color=WHITE)

add_rect(slide, 0.25, 1.2, 4.0, 5.8, GRAY_LIGHT)
add_rect(slide, 0.25, 1.2, 4.0, 0.45, GREEN)
add_text(slide, '대상 회원', 0.4, 1.27, 3.7, 0.32, size=13, bold=True, color=WHITE)

members_info2 = [
    ('정하윤', '수학2A · 국어3A · 한자5A · SSR'),
    ('서아준', '수학A · 국어4A · 한자5A'),
    ('박시환', '수학2A · 국어3A · 한자5A'),
    ('이은서', '수학2A · 국어2A · 한자5A · 호과2'),
]
for i, (name, subj) in enumerate(members_info2):
    y = 1.8 + i * 1.25
    add_rect(slide, 0.35, y, 3.8, 1.1, WHITE)
    add_rect(slide, 0.35, y, 0.08, 1.1, GREEN)
    add_text(slide, name, 0.55, y + 0.06, 3.4, 0.38, size=14, bold=True, color=RGBColor(0x1E, 0x5C, 0x2E))
    add_text(slide, subj, 0.55, y + 0.5, 3.4, 0.45, size=11, color=BLACK)

add_rect(slide, 4.55, 1.2, 8.55, 5.8, GRAY_LIGHT)
add_rect(slide, 4.55, 1.2, 8.55, 0.45, RGBColor(0x1E, 0x5C, 0x2E))
add_text(slide, '상담 화법  (부담 → 효율 전환)', 4.7, 1.27, 8.2, 0.32, size=13, bold=True, color=WHITE)

sections2 = [
    (GREEN, '도입 (공감)',
     '"맞아요, 은서(하윤/시환)처럼 과목이 여러 개면 집에서 다 하기 정말 힘들죠."'),
    (RGBColor(0x1E, 0x5C, 0x2E), '핵심 전환 — "추가가 아닌 효율화"',
     '"스마트구몬N은 과목을 더 추가하는 게 아니에요.\n지금 하고 있는 수학, 국어, 한자에 앱 학습을 연결해주는 방식이에요.\n새로운 과목이 아니라 지금 하는 걸 더 효과적으로 만드는 거예요."'),
    (ORANGE, '앱 동기 설명',
     '"종이 교재 마치면 패드에서 확인 미션이 생겨요.\n아이 입장에서는 패드 하려고 종이 먼저 한다는 동기가 생기는 거예요.\n오히려 지금 학습량을 더 잘 소화하게 도와주는 역할이에요."'),
    (GREEN, '클로징',
     '"한 달만 써보세요. 집에서 하는 게 달라져요."'),
]
for i, (sc, label, text) in enumerate(sections2):
    y = 1.82 + i * 1.25
    add_rect(slide, 4.65, y, 8.35, 1.1, WHITE)
    add_rect(slide, 4.65, y, 0.08, 1.1, sc)
    add_text(slide, label, 4.85, y + 0.04, 7.9, 0.3, size=11, bold=True, color=sc)
    add_text(slide, text, 4.85, y + 0.38, 7.9, 0.65, size=11, color=BLACK)

# ────────────────────────────────────────────────────────────
# 슬라이드 5: 유형 3 — 운필·쓰기 훈련 중심
# ────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
add_rect(slide, 0, 0, 13.33, 1.05, ORANGE)
add_rect(slide, 0, 0, 0.55, 1.05, RGBColor(0xBF, 0x55, 0x00))
add_text(slide, '유형 3', 0.02, 0.28, 0.52, 0.5,
         size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, '운필·쓰기 훈련 중심  |  채아린(유7) · 채아성(유5)', 0.65, 0.18, 12, 0.7,
         size=22, bold=True, color=WHITE)

add_rect(slide, 0.25, 1.2, 4.0, 5.8, GRAY_LIGHT)
add_rect(slide, 0.25, 1.2, 4.0, 0.45, ORANGE)
add_text(slide, '대상 회원', 0.4, 1.27, 3.7, 0.32, size=13, bold=True, color=WHITE)

members_info3 = [
    ('채아린', '유7', '수학2A · 국어3A · 운필4\n(손 훈련 + 기본 학습 병행)'),
    ('채아성', '유5', '국어3A · 한자4A · 운필7\n(5세에 한자까지 — 학습 의욕 높음)'),
]
for i, (name, grade, subj) in enumerate(members_info3):
    y = 1.85 + i * 2.6
    add_rect(slide, 0.35, y, 3.8, 2.3, WHITE)
    add_rect(slide, 0.35, y, 0.08, 2.3, ORANGE)
    add_text(slide, f'{name}  ({grade})', 0.55, y + 0.1, 3.4, 0.45, size=15, bold=True, color=RGBColor(0xBF, 0x55, 0x00))
    add_text(slide, subj, 0.55, y + 0.65, 3.4, 1.4, size=12, color=BLACK)

add_rect(slide, 0.25, 6.35, 4.0, 0.5, ORANGE_LIGHT)
add_text(slide, '★ 4가지 유형 중 앱 시너지 가장 높음', 0.35, 6.4, 3.8, 0.4,
         size=12, bold=True, color=RGBColor(0xBF, 0x55, 0x00))

add_rect(slide, 4.55, 1.2, 8.55, 5.8, GRAY_LIGHT)
add_rect(slide, 4.55, 1.2, 8.55, 0.45, RGBColor(0xBF, 0x55, 0x00))
add_text(slide, '상담 화법  (운필 단계 = 앱 연결 최적 시기)', 4.7, 1.27, 8.2, 0.32, size=13, bold=True, color=WHITE)

sections3 = [
    (ORANGE, '도입 — 운필 단계 의미 설명',
     '"아린이(아성이)가 지금 운필 하고 있잖아요.\n이 단계가 바로 손 훈련이랑 앱 흥미를 동시에 잡을 수 있는 가장 좋은 시기예요."'),
    (RGBColor(0xBF, 0x55, 0x00), '차별점 — 종이와 앱의 역할 분리',
     '"종이로는 쓰는 힘을 키우고, 앱으로는 쓴 내용을 게임처럼 확인해요.\n같은 내용을 다른 방식으로 두 번 경험하니 자연스럽게 기억에 남아요.\n억지로 가르치지 않아도 아이가 스스로 하고 싶어해요."'),
    (GREEN, '화면 걱정 해소',
     '"스마트구몬N 앱은 유튜브나 게임이 아니에요.\n학습 보조 도구라 시간도 짧고, 다른 콘텐츠로 빠져나갈 수가 없어요.\n오히려 패드를 공부 도구로 인식시켜두는 게 나중에 훨씬 도움이 돼요."'),
    (ORANGE, '클로징',
     '"지금 이 시기 놓치면 아까워요. 이번 달 안에 시작해보세요."'),
]
for i, (sc, label, text) in enumerate(sections3):
    y = 1.82 + i * 1.25
    add_rect(slide, 4.65, y, 8.35, 1.1, WHITE)
    add_rect(slide, 4.65, y, 0.08, 1.1, sc)
    add_text(slide, label, 4.85, y + 0.04, 7.9, 0.3, size=11, bold=True, color=sc)
    add_text(slide, text, 4.85, y + 0.38, 7.9, 0.65, size=11, color=BLACK)

# ────────────────────────────────────────────────────────────
# 슬라이드 6: 유형 4 — 국어 단독
# ────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
add_rect(slide, 0, 0, 13.33, 1.05, PURPLE)
add_rect(slide, 0, 0, 0.55, 1.05, RGBColor(0x4A, 0x1A, 0x70))
add_text(slide, '유형 4', 0.02, 0.28, 0.52, 0.5,
         size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, '국어 단독  |  임윤아 (유7 · 국어A)', 0.65, 0.18, 12, 0.7,
         size=22, bold=True, color=WHITE)

add_rect(slide, 0.25, 1.2, 4.0, 5.8, GRAY_LIGHT)
add_rect(slide, 0.25, 1.2, 4.0, 0.45, PURPLE)
add_text(slide, '대상 회원', 0.4, 1.27, 3.7, 0.32, size=13, bold=True, color=WHITE)

add_rect(slide, 0.35, 1.8, 3.8, 1.8, WHITE)
add_rect(slide, 0.35, 1.8, 0.08, 1.8, PURPLE)
add_text(slide, '임윤아  (유7)', 0.55, 1.9, 3.4, 0.45, size=15, bold=True, color=RGBColor(0x4A, 0x1A, 0x70))
add_text(slide, '국어A  (1161~170번)\n현재 국어 단독 학습', 0.55, 2.45, 3.4, 0.9, size=12, color=BLACK)

add_rect(slide, 0.25, 3.85, 4.0, 3.15, PURPLE_LIGHT)
add_rect(slide, 0.25, 3.85, 4.0, 0.4, PURPLE)
add_text(slide, '상담 전략 포인트', 0.4, 3.9, 3.7, 0.3, size=12, bold=True, color=WHITE)
add_text(slide,
         '① 국어 단독 → 수학 추가 연결 기회\n\n'
         '② 앱으로 흥미 붙이면\n    수학도 자연스럽게 시작 가능\n\n'
         '③ 지금 A단계는 시각화 보완이\n    중요한 시기',
         0.4, 4.35, 3.7, 2.5, size=12, color=PURPLE)

add_rect(slide, 4.55, 1.2, 8.55, 5.8, GRAY_LIGHT)
add_rect(slide, 4.55, 1.2, 8.55, 0.45, RGBColor(0x4A, 0x1A, 0x70))
add_text(slide, '상담 화법  (단독 → 확장 전략)', 4.7, 1.27, 8.2, 0.32, size=13, bold=True, color=WHITE)

sections4 = [
    (PURPLE, '도입 — 현재 학습 인정',
     '"윤아가 지금 국어 열심히 하고 있잖아요.\n근데 국어만 하다 보면 수학이랑 같이 하는 친구들이랑 나중에 차이가 날 수 있어요."'),
    (RGBColor(0x4A, 0x1A, 0x70), '전환 — 앱으로 흥미 연결',
     '"스마트구몬N은 종이랑 앱을 같이 써요.\n앱에서 게임처럼 학습하다 보면 수학에도 자연스럽게 관심이 생겨요.\nA단계처럼 기초 개념 시기에 앱 시각화가 엄청 도움이 돼요."'),
    (BLUE_MID, '수학 추가 연결',
     '"처음엔 국어만 앱으로 연결해도 돼요.\n익숙해지면 수학도 같이 시작해보는 거예요.\n지금 시작해두면 초등 올라갈 때 완전히 다른 출발선에 서요."'),
    (PURPLE, '클로징',
     '"이번 달 신청하면 제일 좋은 조건이에요. 한번 같이 해봐요."'),
]
for i, (sc, label, text) in enumerate(sections4):
    y = 1.82 + i * 1.25
    add_rect(slide, 4.65, y, 8.35, 1.1, WHITE)
    add_rect(slide, 4.65, y, 0.08, 1.1, sc)
    add_text(slide, label, 4.85, y + 0.04, 7.9, 0.3, size=11, bold=True, color=sc)
    add_text(slide, text, 4.85, y + 0.38, 7.9, 0.65, size=11, color=BLACK)

# ────────────────────────────────────────────────────────────
# 슬라이드 7: 가격 + 인센티브 클로징
# ────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
add_rect(slide, 0, 0, 13.33, 1.05, RED)
add_text(slide, '이번 달 신청이 제일 좋은 조건입니다', 0.3, 0.18, 12.5, 0.7,
         size=26, bold=True, color=WHITE)

# 3개 박스
boxes = [
    (BLUE_MID,  '월 추가 비용', '+ 23,000원', '기존 구몬 비용에\n월 23,000원만 추가'),
    (ORANGE,    'N페이 지급', '1만원 캐시백', '이번 달 신청 시\nN페이 10,000원 지급'),
    (RED,       '다음 달부터', '25,000원 인상', '이번 달 신청하면\n현재 요금 유지 가능'),
]
for i, (bc, label, main, sub) in enumerate(boxes):
    x = 0.5 + i * 4.2
    add_rect(slide, x, 1.3, 3.8, 3.8, bc)
    add_text(slide, label, x, 1.45, 3.8, 0.5,
             size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, main, x, 2.1, 3.8, 0.9,
             size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, x + 0.2, 3.1, 3.4, 0.04, WHITE)
    add_text(slide, sub, x, 3.25, 3.8, 0.9,
             size=14, color=WHITE, align=PP_ALIGN.CENTER)

# 계산
add_rect(slide, 0.5, 5.35, 12.33, 1.4, RGBColor(0xFD, 0xF2, 0xF2))
add_rect(slide, 0.5, 5.35, 0.08, 1.4, RED)
add_text(slide, '실질 첫 달 계산  :  23,000원 − N페이 10,000원 = 실부담 13,000원  |  하루 환산 750원',
         0.7, 5.55, 12.0, 0.5, size=16, bold=True, color=RED)
add_text(slide, '※ 다음 달부터 25,000원으로 인상 예정  |  이번 달 신청 시 현재 요금 유지',
         0.7, 6.1, 12.0, 0.4, size=13, color=RGBColor(0x80, 0x00, 0x00))

out = 'C:/Users/najun/Downloads/스마트구몬N_유형별_설득화법.pptx'
prs.save(out)
print(f'완료: {os.path.getsize(out):,} bytes  ({os.path.getsize(out)/1024:.0f} KB)')
print(f'저장 위치: {out}')
