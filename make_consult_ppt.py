import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

DESKTOP = 'C:/Users/najun/OneDrive/Desktop/'

def add_image(slide, path, l, t, w, h):
    try:
        slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))
    except Exception as e:
        print(f'이미지 삽입 실패: {path} - {e}')

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── 색상 ────────────────────────────────────────────────
KUMON_RED  = RGBColor(0xE8, 0x4B, 0x2A)   # 구몬 오렌지레드
DARK       = RGBColor(0x1F, 0x35, 0x64)
BLUE       = RGBColor(0x2E, 0x75, 0xB6)
ORANGE     = RGBColor(0xED, 0x7D, 0x31)
GREEN      = RGBColor(0x70, 0xAD, 0x47)
RED        = RGBColor(0xC0, 0x00, 0x00)
PURPLE     = RGBColor(0x70, 0x30, 0xA0)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY      = RGBColor(0xF5, 0xF5, 0xF5)
YELLOW_L   = RGBColor(0xFF, 0xF2, 0xCC)
BLACK      = RGBColor(0x00, 0x00, 0x00)

blank = prs.slide_layouts[6]

def rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.line.fill.background()
    s.fill.solid()
    s.fill.fore_color.rgb = color
    return s

def txt(slide, text, l, t, w, h, size=13, bold=False,
        color=None, align=PP_ALIGN.LEFT, wrap=True):
    if color is None: color = BLACK
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb

def multi_txt(slide, lines, l, t, w, h, size=13, bold=False, color=None,
              align=PP_ALIGN.LEFT, line_spacing=1.15):
    if color is None: color = BLACK
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    from pptx.util import Pt as PT
    from pptx.oxml.ns import qn
    import lxml.etree as etree
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(2)
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb

def badge(slide, text, l, t, w, h, bg, fg=None):
    if fg is None: fg = WHITE
    rect(slide, l, t, w, h, bg)
    txt(slide, text, l+0.05, t+0.04, w-0.1, h-0.08,
        size=11, bold=True, color=fg, align=PP_ALIGN.CENTER)

def stat_card(slide, label, value, unit, l, t, w, h, bg):
    rect(slide, l, t, w, h, bg)
    rect(slide, l, t, w, 0.45, bg)
    txt(slide, label, l+0.1, t+0.06, w-0.2, 0.35,
        size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(slide, value, l+0.1, t+0.5, w-0.2, 0.7,
        size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(slide, unit, l+0.1, t+1.2, w-0.2, 0.3,
        size=10, color=WHITE, align=PP_ALIGN.CENTER)

def speech_bubble(slide, lines, l, t, w, h, bg=None, fg=None):
    if bg is None: bg = YELLOW_L
    if fg is None: fg = DARK
    rect(slide, l, t, w, h, bg)
    multi_txt(slide, lines, l+0.2, t+0.15, w-0.35, h-0.3,
              size=12, color=fg)

# ════════════════════════════════════════════════════════
# 슬라이드 1: 표지
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
rect(slide, 0, 0, 13.33, 7.5, DARK)
rect(slide, 0, 5.5, 13.33, 2.0, KUMON_RED)
rect(slide, 0, 3.3, 13.33, 0.06, WHITE)

txt(slide, '마스타 입회 상담', 0, 1.0, 13.33, 1.2,
    size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(slide, '장학금 연계 맞춤 상담 카드', 0, 2.3, 13.33, 0.8,
    size=24, color=RGBColor(0xBD,0xD7,0xEE), align=PP_ALIGN.CENTER)

txt(slide, '정유찬  ·  강윤서', 0, 3.6, 13.33, 0.8,
    size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(slide, '구몬 관양3지구  |  박양희 선생님  ·  김은숙 선생님', 0, 4.3, 13.33, 0.6,
    size=15, color=RGBColor(0xFF,0xC7,0xCE), align=PP_ALIGN.CENTER)
txt(slide, '2026년 4월  |  마스타 입회 상담 자료', 0, 6.0, 13.33, 0.5,
    size=13, color=WHITE, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════
# 슬라이드 2: 정유찬 — 현황 대시보드
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
rect(slide, 0, 0, 13.33, 1.2, KUMON_RED)
txt(slide, '정유찬  |  초3 · 수학  |  현황 대시보드', 0.3, 0.2, 9, 0.8,
    size=26, bold=True, color=WHITE)
txt(slide, '박양희 선생님', 10.0, 0.3, 3.0, 0.5,
    size=14, color=RGBColor(0xFF,0xC7,0xCE), align=PP_ALIGN.RIGHT)

# 스탯 카드 4개
stat_card(slide, '전국 진도순위', '상위 3%', '초3 전체 기준', 0.3, 1.4, 2.9, 1.7, KUMON_RED)
stat_card(slide, '현재 진도',    'F단계',   '초6 수준',       3.4, 1.4, 2.9, 1.7, BLUE)
stat_card(slide, '학년 대비 갭', '+3단계',  '가산점 +30점',   6.5, 1.4, 2.9, 1.7, GREEN)
stat_card(slide, '최종대회',     'D합격',   '학습기간 54개월', 9.6, 1.4, 2.9, 1.7, PURPLE)

# 정유찬 이미지 삽입 (순위 + 그래프)
add_image(slide, DESKTOP + '정유찬 순위.jpg',   0.3, 3.25, 4.0, 4.0)
add_image(slide, DESKTOP + '정유찬 글래프.jpg', 4.5, 3.25, 4.1, 4.0)

# 장학금 로드맵
rect(slide, 9.0, 3.3, 4.1, 3.9, LGRAY)
rect(slide, 9.0, 3.3, 4.1, 0.45, KUMON_RED)
txt(slide, '장학금 로드맵', 9.2, 3.35, 3.7, 0.38,
    size=13, bold=True, color=WHITE)

roadmap = [
    ('현재', 'F단계', '초6수준', DARK),
    ('→', '', '', WHITE),
    ('주니어', '상위 10명', '상장+전집', BLUE),
    ('→', '', '', WHITE),
    ('N단계', '합격 시', '300만원', GREEN),
    ('→', '', '', WHITE),
    ('O단계', '합격 시', '400만원', KUMON_RED),
]
y = 3.9
for label, sub1, sub2, color in roadmap:
    if label == '→':
        txt(slide, '▼', 10.8, y, 0.5, 0.3, size=12, color=DARK, align=PP_ALIGN.CENTER)
        y += 0.28
    else:
        rect(slide, 9.2, y, 3.7, 0.6, color)
        txt(slide, label, 9.3, y+0.05, 1.2, 0.5, size=11, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)
        txt(slide, sub1, 10.55, y+0.05, 1.3, 0.25, size=10,
            color=WHITE, align=PP_ALIGN.CENTER)
        txt(slide, sub2, 10.55, y+0.3, 1.3, 0.25, size=10,
            color=YELLOW_L, align=PP_ALIGN.CENTER)
        y += 0.68

# ════════════════════════════════════════════════════════
# 슬라이드 3: 정유찬 — 상담 멘트 카드
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
rect(slide, 0, 0, 13.33, 1.1, KUMON_RED)
txt(slide, '정유찬  |  마스타 입회 상담 멘트', 0.3, 0.15, 12, 0.8,
    size=24, bold=True, color=WHITE)

steps = [
    (KUMON_RED, 'OPEN',  '전국 상위 3%로 시작',
     '"어머니, 정유찬이 지금 전국 구몬 초3 회원 중 상위 3%예요.\n100명 중 3등 안에 드는 거거든요."'),
    (BLUE,      'FACT',  '그래프로 시각화',
     '"초3인데 초6 수준(F단계)을 배우고 있어요. 그리고 이 격차가\n시간이 갈수록 점점 더 벌어지고 있어요. (그래프 보여주며)"'),
    (GREEN,     'AWARD', '두 가지 시상 기회',
     '"① 주니어장학금 — 학년별 상위 10명, 전국 3%면 충분해요.\n② 대학장학금 — N단계 합격 시 300만 원, O단계 합격 시 400만 원"'),
    (PURPLE,    'NEED',  '가산점 구조 설명',
     '"3단계 앞서 있으니 지금 +30점 가산점이 쌓이고 있어요.\n이 속도가 떨어지면 가산점도 줄어드니 지금이 골든타임이에요."'),
    (DARK,      'CLOSE', '마스타 연결',
     '"초3 때 마스타 들어온 아이들이 나중에 대학장학금 받은\n케이스가 있어요. 다음 달 자리 한번 넣어보시겠어요?"'),
]

for i, (color, step, title, script) in enumerate(steps):
    col = i % 3
    row = i // 3
    x = 0.3 + col * 4.35
    y = 1.25 + row * 2.9

    rect(slide, x, y, 4.1, 2.6, LGRAY)
    rect(slide, x, y, 4.1, 0.5, color)
    badge(slide, step, x+0.1, y+0.08, 0.85, 0.35, WHITE, color)
    txt(slide, title, x+1.05, y+0.1, 2.9, 0.35,
        size=12, bold=True, color=WHITE)
    txt(slide, script, x+0.15, y+0.62, 3.8, 1.85,
        size=11, color=DARK)

# 하단 핵심 메시지
rect(slide, 0, 7.0, 13.33, 0.5, YELLOW_L)
txt(slide, '💡  핵심 메시지: "초3에 전국 3% — 이 속도를 마스타로 연결하면 대학장학금이 현실이 됩니다"',
    0.4, 7.05, 12.5, 0.4, size=13, bold=True, color=DARK)

# ════════════════════════════════════════════════════════
# 슬라이드 4: 강윤서 — 현황 대시보드
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
rect(slide, 0, 0, 13.33, 1.2, DARK)
txt(slide, '강윤서  |  중2 · 수플  |  현황 대시보드', 0.3, 0.2, 9, 0.8,
    size=26, bold=True, color=WHITE)
txt(slide, '김은숙 선생님', 10.0, 0.3, 3.0, 0.5,
    size=14, color=RGBColor(0x9D,0xC3,0xE6), align=PP_ALIGN.RIGHT)

stat_card(slide, '전국 진도순위', '상위 40%', '중2 전체 기준', 0.3, 1.4, 2.9, 1.7, KUMON_RED)
stat_card(slide, '현재 진도',    'K단계',   '고2 수준',       3.4, 1.4, 2.9, 1.7, DARK)
stat_card(slide, '학년 대비 갭', '+3단계',  '가산점 +30점',   6.5, 1.4, 2.9, 1.7, GREEN)
stat_card(slide, '최종대회',     'I합격',   '학습기간 15개월', 9.6, 1.4, 2.9, 1.7, BLUE)

# 강윤서 이미지 삽입 (순위 + 그래프)
add_image(slide, DESKTOP + '강윤서 순위.jpg',   0.3, 3.25, 4.0, 4.0)
add_image(slide, DESKTOP + '강윤서 그래프.jpg', 4.5, 3.25, 4.1, 4.0)

# 장학금 로드맵
rect(slide, 9.0, 3.3, 4.1, 3.9, LGRAY)
rect(slide, 9.0, 3.3, 4.1, 0.45, DARK)
txt(slide, '대학장학금 로드맵', 9.2, 3.35, 3.7, 0.38,
    size=13, bold=True, color=WHITE)
multi_txt(slide, [
    '현재: K단계 (고2수준)',
    '',
    '▼  L단계 (고3수준)',
    '',
    '▼  M단계 (수능수준)',
    '',
    '▼  N단계 합격  →  300만 원',
    '',
    '▼  O단계 합격  →  400만 원',
    '',
    '★  K → N까지 딱 3단계!',
], 9.3, 3.85, 3.6, 3.2, size=12, color=DARK)

# N/O 강조
rect(slide, 9.2, 6.3, 3.7, 0.65, GREEN)
txt(slide, 'K→N: 3단계  |  예상 수령액 300만 원+', 9.3, 6.38, 3.5, 0.5,
    size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════
# 슬라이드 5: 강윤서 — 상담 멘트 카드
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
rect(slide, 0, 0, 13.33, 1.1, DARK)
txt(slide, '강윤서  |  마스타 입회 상담 멘트', 0.3, 0.15, 12, 0.8,
    size=24, bold=True, color=WHITE)

steps2 = [
    (DARK,    'OPEN',  '진도 현황으로 시작',
     '"어머니, 오늘 꼭 드리고 싶은 말씀이 있어서요.\n강윤서가 지금 중2인데 고2 수준(K단계)을 배우고 있어요."'),
    (BLUE,    'FACT',  '장학금까지 거리 제시',
     '"N단계 합격하면 300만 원, O단계 합격하면 400만 원이에요.\nK→N까지 딱 3단계밖에 안 남았어요. (손가락으로 짚으며)"'),
    (GREEN,   'AWARD', '가산점 강조',
     '"3단계 앞서 있으니 지금 +30점 가산점이 붙고 있어요.\n최종대회 I단계 합격 이력도 있으니 수상 선정에 유리해요."'),
    (ORANGE,  'NEED',  '마스타 필요성',
     '"이 진도라면 일반 수업 속도가 느릴 수 있어요.\n마스타는 상위 진도 회원을 위한 심화 과정이에요."'),
    (KUMON_RED,'CLOSE', '비용 역전 프레이밍',
     '"수강료 내면서 열심히 하다가 300만 원 돌려받는 구조예요.\n다음 달 마스타 자리 한번 넣어드릴까요?"'),
]

for i, (color, step, title, script) in enumerate(steps2):
    col = i % 3
    row = i // 3
    x = 0.3 + col * 4.35
    y = 1.25 + row * 2.9

    rect(slide, x, y, 4.1, 2.6, LGRAY)
    rect(slide, x, y, 4.1, 0.5, color)
    badge(slide, step, x+0.1, y+0.08, 0.85, 0.35, WHITE, color)
    txt(slide, title, x+1.05, y+0.1, 2.9, 0.35,
        size=12, bold=True, color=WHITE)
    txt(slide, script, x+0.15, y+0.62, 3.8, 1.85,
        size=11, color=DARK)

rect(slide, 0, 7.0, 13.33, 0.5, YELLOW_L)
txt(slide, '💡  핵심 메시지: "K단계에서 N단계까지 딱 3단계 — 300만 원 장학금이 현실입니다"',
    0.4, 7.05, 12.5, 0.4, size=13, bold=True, color=DARK)

# ── 저장 ───────────────────────────────────────────────
out = 'C:/Users/najun/Downloads/마스타_상담카드_정유찬_강윤서.pptx'
prs.save(out)
import os
print(f'완료: {out}')
print(f'크기: {os.path.getsize(out):,} bytes ({os.path.getsize(out)/1024:.0f} KB)')
