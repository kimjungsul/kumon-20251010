import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

BLUE_DARK = RGBColor(0x1F, 0x35, 0x64)
BLUE_MID = RGBColor(0x2E, 0x75, 0xB6)
BLUE_LIGHT = RGBColor(0xD6, 0xE4, 0xF0)
ORANGE = RGBColor(0xED, 0x7D, 0x31)
GREEN = RGBColor(0x70, 0xAD, 0x47)
RED = RGBColor(0xC0, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
GRAY = RGBColor(0xF2, 0xF2, 0xF2)


def add_rect(slide, l, t, w, h, color):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    return shape


def add_text(slide, text, l, t, w, h, size=18, bold=False, color=None, align=PP_ALIGN.LEFT):
    if color is None:
        color = BLACK
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def add_bullets(slide, lines, l, t, w, h, size=18, color=None):
    if color is None:
        color = BLACK
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.clear()
    first = True
    for text, level in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = text
        p.level = level
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = False
        p.font.name = 'Malgun Gothic'
        p.space_before = Pt(3)
        p.space_after = Pt(0)
        p.line_spacing = 1.2
        p.font.name = 'Malgun Gothic'
    return txBox

# Slide 1: Title
slide = prs.slides.add_slide(blank)
add_rect(slide, 0, 0, 13.33, 7.5, BLUE_DARK)
add_text(slide, '어린이날 연휴, 스마트한 교재 관리 매뉴얼', 0.5, 1.3, 12.3, 1.2,
         size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, '진도 설정부터 효율적인 채점 피드백까지', 0.5, 2.6, 12.3, 0.8,
         size=24, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, '대상: 1년 차 미만 선생님을 위한 실전 가이드', 0.5, 4.4, 12.3, 0.6,
         size=18, color=WHITE, align=PP_ALIGN.CENTER)

# Slide 2: 핵심 목적
slide = prs.slides.add_slide(blank)
add_rect(slide, 0, 0, 13.33, 1.0, BLUE_MID)
add_text(slide, '이번 교육의 핵심 목적', 0.4, 0.16, 12.5, 0.7,
         size=28, bold=True, color=WHITE)
add_bullets(slide, [
    ('회원 부담 경감: 연휴 기간 학습 거부감 최소화', 0),
    ('학습 연속성 유지: 진도 조절을 통한 \'밀린 교재\' 방지', 0),
    ('업무 효율화: 2주분 교재의 효과적인 채점 및 관리', 0),
], 0.8, 1.2, 11.7, 5.5, size=22, color=BLACK)

# Slide 3: 진도 설정 전략
slide = prs.slides.add_slide(blank)
add_rect(slide, 0, 0, 13.33, 1.0, ORANGE)
add_text(slide, '핵심 진도 설정 전략 - "복습이 답이다"', 0.4, 0.16, 12.5, 0.75,
         size=28, bold=True, color=WHITE)
add_bullets(slide, [
    ('원칙: 새로운 개념보다는 익숙한 복습 위주로 구성', 0),
    ('이유: 어려운 새 단계가 연휴와 겹치면 아이의 학습 의욕이 꺾이고 퇴회 위험이 높아짐', 0),
    ('방법: 아이가 혼자서도 5~10분 내에 끝낼 수 있는 수준으로 배정', 0),
    ('"연휴 동안은 가볍게 복습하고, 선생님이랑 다음 수업 때 기분 좋게 새 진도 나가자!"라고 동기부여', 0),
], 0.8, 1.2, 11.7, 5.5, size=20, color=BLACK)

# Slide 4: 체크리스트
slide = prs.slides.add_slide(blank)
add_rect(slide, 0, 0, 13.33, 1.0, GREEN)
add_text(slide, '교재 전달 전 체크리스트', 0.4, 0.16, 12.5, 0.75,
         size=28, bold=True, color=WHITE)
add_bullets(slide, [
    ('준비: 1주 차/2주 차 분량 명확히 구분 (인덱스 활용)', 0),
    ('포장: 어린이날 축하 메시지나 작은 간식 동봉 (선물 느낌 강조)', 0),
    ('예고: 학부모에게 미리 휴강 및 교재 전달 사항 안내', 0),
], 0.8, 1.2, 11.7, 5.5, size=20, color=BLACK)

# Slide 5: 매직 스크립트
slide = prs.slides.add_slide(blank)
add_rect(slide, 0, 0, 13.33, 1.0, BLUE_DARK)
add_text(slide, '상황별 바로 쓰는 매직 스크립트', 0.4, 0.16, 12.5, 0.75,
         size=28, bold=True, color=WHITE)
add_bullets(slide, [
    ('대면 시: "어머님, 이번에는 연휴 동안 부담 없이 할 수 있는 복습 위주 교재로 준비해 두었습니다. 쉬운 분량이라 아이도 편하게 풀 수 있을 거예요."', 0),
    ('메시지형: "어린이날 잘 보내시길 바라요. 이번 교재는 아이가 혼자서도 부담 없이 풀 수 있는 복습 위주 자료로 준비했습니다. 가방에서 확인 부탁드릴게요."', 0),
], 0.8, 1.2, 11.7, 5.5, size=20, color=BLACK)

# Slide 6: 채점 및 피드백 노하우
slide = prs.slides.add_slide(blank)
add_rect(slide, 0, 0, 13.33, 1.0, RED)
add_text(slide, '2주분 교재 채점 및 피드백 노하우', 0.4, 0.16, 12.5, 0.75,
         size=28, bold=True, color=WHITE)
add_bullets(slide, [
    ('중복 내용 처리: 내용이 중복되는 부분은 한 세트만 집중 채점하여 선생님의 업무 과부하 방지', 0),
    ('오답 위주 관리: 전체를 다 보기보다 틀린 부분 위주로 정밀 피드백 실시', 0),
    ('지침:', 0),
    ('채점용 교재를 반드시 지참하여 현장에서 즉시 피드백', 1),
    ('아이가 반복해서 틀리는 지점만 명확히 짚어주어 수업 시간 효율 극대화', 1),
], 0.8, 1.2, 11.7, 5.5, size=20, color=BLACK)

# Slide 7: 실전 Tip
slide = prs.slides.add_slide(blank)
add_rect(slide, 0, 0, 13.33, 1.0, BLUE_MID)
add_text(slide, '신입 선생님을 위한 실전 Tip', 0.4, 0.16, 12.5, 0.75,
         size=28, bold=True, color=WHITE)
add_bullets(slide, [
    ('자신감: 연휴도 아이의 학습 흐름을 지켜주는 중요한 기간이라는 마음가짐 갖기', 0),
    ('메모 활용: 교재 첫 장에 "연휴 잘 보내고 다음 주에 웃으며 만나자!"라는 응원 문구 남기기', 0),
], 0.8, 1.2, 11.7, 5.5, size=20, color=BLACK)

# Slide 8: 마무리
slide = prs.slides.add_slide(blank)
add_rect(slide, 0, 0, 13.33, 1.0, GREEN)
add_text(slide, '마무리 및 Q&A', 0.4, 0.16, 12.5, 0.75,
         size=28, bold=True, color=WHITE)
add_text(slide, '세심한 진도 조절과 정확한 피드백이 \'진짜 실력\'입니다.',
         0.8, 1.4, 11.7, 2.5, size=24, color=BLACK)

out = os.path.join(os.path.dirname(__file__), '어린이날_연휴_교재_관리_매뉴얼.pptx')
prs.save(out)
print(f'완료: {out}')
