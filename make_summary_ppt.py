import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLUE_DARK  = RGBColor(0x1F, 0x35, 0x64)
BLUE_MID   = RGBColor(0x2E, 0x75, 0xB6)
BLUE_LIGHT = RGBColor(0xD6, 0xE4, 0xF0)
RED        = RGBColor(0xC0, 0x00, 0x00)
ORANGE     = RGBColor(0xED, 0x7D, 0x31)
GREEN      = RGBColor(0x70, 0xAD, 0x47)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_LIGHT = RGBColor(0xF2, 0xF2, 0xF2)
BLACK      = RGBColor(0x00, 0x00, 0x00)

blank = prs.slide_layouts[6]

def add_rect(slide, l, t, w, h, color):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    return shape

def add_text(slide, text, l, t, w, h, size=18, bold=False, color=None, align=PP_ALIGN.LEFT):
    if color is None:
        color = BLACK
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox

# ── 슬라이드 1: 표지 ──────────────────────────────────────
slide = prs.slides.add_slide(blank)
add_rect(slide, 0, 0, 13.33, 7.5, BLUE_DARK)
add_rect(slide, 0, 5.8, 13.33, 1.7, BLUE_MID)
add_rect(slide, 1.5, 1.0, 10.33, 0.08, WHITE)

add_text(slide, '한글 기본 지도법', 0, 1.5, 13.33, 1.5,
         size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, '5편 · 6편  핵심 요약', 0, 3.0, 13.33, 0.8,
         size=28, color=RGBColor(0xBD, 0xD7, 0xEE), align=PP_ALIGN.CENTER)
add_text(slide, '나무단계 1블록 · 2블록 지도 포인트', 0, 3.8, 13.33, 0.6,
         size=18, color=RGBColor(0x9D, 0xC3, 0xE6), align=PP_ALIGN.CENTER)
add_text(slide, '구몬 관양지국  |  교사 교육 자료', 0, 6.3, 13.33, 0.5,
         size=14, color=WHITE, align=PP_ALIGN.CENTER)

# ── 슬라이드 2: 핵심 원칙 ─────────────────────────────────
slide = prs.slides.add_slide(blank)
add_rect(slide, 0, 0, 13.33, 1.1, BLUE_DARK)
add_text(slide, '영유아 지도 핵심 원칙 (5편·6편 공통)', 0.3, 0.15, 12, 0.8,
         size=24, bold=True, color=WHITE)

cards = [
    ('즐겁게 학습한다', 'BLUE_MID',
     '유치원 다니는 시기까지 즐겁게 학습할 수 있도록 지도\n\n도전하는 마음을 가질 수 있도록 이끌어 준다\n\n정답을 강요하지 않고 긍정적인 분위기 유지'),
    ('자립을 위한 지도', 'GREEN',
     '스스로 할 수 있는 힘을 기르는 것을 최종 목표\n\n정답을 알려주기보다 스스로 생각하게 유도\n\n나무단계 = 새싹단계와 기본 지도방법 동일'),
]
colors = [BLUE_MID, GREEN]

for i, (title, _, body) in enumerate(cards):
    x = 0.5 + i * 6.4
    add_rect(slide, x, 1.3, 6.0, 5.5, GRAY_LIGHT)
    add_rect(slide, x, 1.3, 6.0, 1.0, colors[i])
    add_text(slide, '★  ' + title, x + 0.2, 1.38, 5.6, 0.8,
             size=20, bold=True, color=WHITE)
    add_text(slide, body, x + 0.3, 2.5, 5.5, 4.0,
             size=16, color=BLACK)

add_rect(slide, 0, 6.85, 13.33, 0.65, BLUE_LIGHT)
add_text(slide, '핵심: 블록별 특징과 지도 포인트 중심으로 운영 — 교재 지도법은 새싹단계와 동일',
         0.4, 6.9, 12.5, 0.55, size=14, bold=True, color=BLUE_DARK)

# ── 슬라이드 3: 5편 — 통글자 지도법 ─────────────────────────
slide = prs.slides.add_slide(blank)
add_rect(slide, 0, 0, 13.33, 1.1, RED)
add_text(slide, '5편 | 통글자 지도법  (나무 1블록)', 0.3, 0.15, 12, 0.8,
         size=24, bold=True, color=WHITE)

points = [
    ('Point 1', '통글자부터 시작하는 이유',
     '유아는 "의미"가 있는 것에만 흥미를 가짐\n의미의 최소 단위 = 낱말  →  통글자부터 노출'),
    ('Point 2', '핵심 블록은 3블록!',
     '1블록: 통글자 익히기 출발점  |  3블록: 한글 읽기 완성의 핵심\n1블록 과제단어 96개 (이 중 37개는 새싹단계 연계어휘)'),
    ('Point 3', '100% 암기는 NO!',
     '통글자만으로 한글 완성 시 11,172자 필요 (현실 불가)\n아이가 관심 있는 2~3개 단어 집중 지도로도 충분'),
]

for i, (badge, title, body) in enumerate(points):
    y = 1.25 + i * 1.95
    add_rect(slide, 0.4, y, 12.5, 1.75, GRAY_LIGHT)
    add_rect(slide, 0.4, y, 1.5, 1.75, RED)
    add_text(slide, badge, 0.42, y + 0.55, 1.46, 0.6,
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, title, 2.1, y + 0.1, 10.5, 0.55,
             size=16, bold=True, color=RED)
    add_text(slide, body, 2.1, y + 0.65, 10.5, 1.0,
             size=13, color=BLACK)

# ── 슬라이드 4: 6편 — 낱글자 지도법 ─────────────────────────
slide = prs.slides.add_slide(blank)
add_rect(slide, 0, 0, 13.33, 1.1, ORANGE)
add_text(slide, '6편 | 낱글자 지도법  (나무 2블록)', 0.3, 0.15, 12, 0.8,
         size=24, bold=True, color=WHITE)

# 좌측 패널
add_rect(slide, 0.3, 1.2, 6.2, 5.8, GRAY_LIGHT)
add_rect(slide, 0.3, 1.2, 6.2, 0.55, ORANGE)
add_text(slide, '낱글자 학습 원리', 0.5, 1.27, 5.8, 0.4,
         size=15, bold=True, color=WHITE)

left_items = [
    ('Point 1', '통글자 → 낱글자 자연스러운 흐름\n통글자 학습 중 "이건 무슨 글자지?" 궁금증 발생\n2블록이 그 궁금증에 답해주는 블록'),
    ('Point 2', '의미 있는 대상으로 낱글자 학습\n예) "가방"의 "가"  →  의미를 가진 단어로 암기\n발음만 가르치면 흥미 유발 어려움'),
]
for i, (pt, txt) in enumerate(left_items):
    y = 1.9 + i * 2.4
    add_rect(slide, 0.5, y, 1.2, 0.45, ORANGE)
    add_text(slide, pt, 0.52, y + 0.05, 1.16, 0.35,
             size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, txt, 0.5, y + 0.6, 5.8, 1.65, size=13, color=BLACK)

# 우측 패널
add_rect(slide, 6.8, 1.2, 6.2, 5.8, GRAY_LIGHT)
add_rect(slide, 6.8, 1.2, 6.2, 0.55, GREEN)
add_text(slide, '낱글자 지도 TIP', 7.0, 1.27, 5.8, 0.4,
         size=15, bold=True, color=WHITE)

right_items = [
    ('TIP 1', '손뼉 치며 낱말 나누기\n손뼉 한 번에 낱자 하나 대응\n리듬감 있게 반복 연습'),
    ('TIP 2', '글자 분리 → 합성 활동\n글자 나누어 보는 연습 진행\n글자가 "떨어지는" 개념 인식'),
]
for i, (tip, txt) in enumerate(right_items):
    y = 1.9 + i * 2.4
    add_rect(slide, 7.0, y, 1.2, 0.45, GREEN)
    add_text(slide, tip, 7.02, y + 0.05, 1.16, 0.35,
             size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, txt, 7.0, y + 0.6, 5.8, 1.65, size=13, color=BLACK)

# ── 슬라이드 5: 그림동화 & 그림카드 지도법 ───────────────────
slide = prs.slides.add_slide(blank)
add_rect(slide, 0, 0, 13.33, 1.1, GREEN)
add_text(slide, '교재 지도법 | 그림동화 & 그림카드 (2블록)', 0.3, 0.15, 12, 0.8,
         size=24, bold=True, color=WHITE)

# 그림동화
add_rect(slide, 0.3, 1.2, 6.1, 5.8, GRAY_LIGHT)
add_rect(slide, 0.3, 1.2, 6.1, 0.55, BLUE_MID)
add_text(slide, '그림동화 지도', 0.5, 1.27, 5.7, 0.4,
         size=15, bold=True, color=WHITE)
add_text(slide, '구몬 그림동화 vs 타사 이야기동화', 0.5, 1.9, 5.7, 0.45,
         size=13, bold=True, color=BLUE_MID)
add_text(slide,
         '• 타사: 그림 속 글 혼합 → 그림 집중 어려움\n'
         '• 구몬: 글·그림 명확히 분리 → 그림 집중 가능\n\n'
         '지도 2 Step\n\n'
         'Step 1.  일러스트 보여주며 동화 전체 읽어주기\n'
         '           (스마트펜·육성 모두 OK / 아이 표정 관찰)\n\n'
         'Step 2.  내용 흐름 파악 후 과제 낱말 인지 확인\n'
         '           (손가락으로 짚어보게끔 지도)',
         0.5, 2.5, 5.7, 4.3, size=13, color=BLACK)

# 그림카드
add_rect(slide, 6.9, 1.2, 6.1, 5.8, GRAY_LIGHT)
add_rect(slide, 6.9, 1.2, 6.1, 0.55, ORANGE)
add_text(slide, '그림카드 지도', 7.1, 1.27, 5.7, 0.4,
         size=15, bold=True, color=WHITE)
add_text(slide, '그림카드 특징', 7.1, 1.9, 5.7, 0.45,
         size=13, bold=True, color=ORANGE)
add_text(slide,
         '특징 1.  글자 수가 다른 낱말 구성\n'
         '           글자 몰라도 낱말 구별 가능\n\n'
         '특징 2.  낱말 구성 글자색이 다름\n'
         '           유아의 색 인지 발달 특성 반영\n\n'
         '지도 방법\n\n'
         '• 플래쉬 기법: 뒤에서 앞으로 보여주기\n'
         '• 1초에 1장씩 스피드 있게 제시\n'
         '• 아이 반응 관찰하며 진행',
         7.1, 2.5, 5.7, 4.3, size=13, color=BLACK)

out = 'd:/coding/kumon-20251010/한글_기본_지도법_5_6편_요약.pptx'
prs.save(out)
print(f'완료: {os.path.getsize(out):,} bytes ({os.path.getsize(out)/1024:.0f} KB)')
